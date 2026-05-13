"""Generate test documents for manual end-to-end testing.

Usage:
    cd AI_Assistant
    python scripts/generate_test_data.py

Produces:
    data/documents/test_report.pdf       ← text-based PDF (for pymupdf)
    data/documents/test_manual.docx      ← Word document
    data/documents/test_slides.pptx      ← PowerPoint
"""

from pathlib import Path

OUT = Path("data/documents")
OUT.mkdir(parents=True, exist_ok=True)


def generate_pdf():
    """Create a text-based PDF with rich content for RAG testing."""
    import fitz

    doc = fitz.open()
    page = doc.new_page()

    y = 72
    title_font = 20
    body_font = 12

    page.insert_text((72, y), "2024年度AI技术发展报告", fontsize=title_font)
    y += 40

    lines = [
        ("一、概述", 16),
        ("", 12),
        ("本报告总结了2024年度人工智能领域的关键技术进展和产业应用趋势。", 12),
        ("报告涵盖大语言模型、多模态AI、AI Agent、开源生态四个主要方向。", 12),
        ("", 12),
        ("二、大语言模型", 16),
        ("", 12),
        ("1. DeepSeek V4于2025年3月发布，采用Mixture-of-Experts架构，总参数量超过1万亿，", 12),
        ("   但每个token仅激活约37亿参数，推理成本大幅降低。其在数学推理、代码生成等", 12),
        ("   基准测试中表现与GPT-4 Turbo相当，API定价仅为后者的1/50。", 12),
        ("", 12),
        ("2. Qwen3系列由阿里通义实验室推出，包含0.6B到72B多个规模。其最大亮点是", 12),
        ("   支持thinking模式的开关切换，用户可根据任务复杂度灵活选择推理深度。", 12),
        ("", 12),
        ("三、多模态AI", 16),
        ("", 12),
        ("1. GPT-4o实现了文本、图像、音频的原生多模态理解，延迟降至232ms。", 12),
        ("2. Claude 4 Opus在视觉推理和长文档分析方面表现突出，支持1M token上下文。", 12),
        ("", 12),
        ("四、AI Agent", 16),
        ("", 12),
        ("1. AI Agent从概念验证走向生产部署，主要框架包括LangGraph、CrewAI和AutoGen。", 12),
        ("2. Anthropic推出的MCP(Model Context Protocol)成为Agent与外部工具交互的标准协议。", 12),
        ("3. Google的A2A(Agent-to-Agent)协议解决了多Agent协作中的通信和发现难题。", 12),
        ("", 12),
        ("五、开源生态", 16),
        ("", 12),
        ("1. LlamaIndex在RAG领域持续领先，提供从文档解析到检索生成的完整Pipeline。", 12),
        ("2. LangChain向LangGraph转型，聚焦复杂Agent工作流的编排和状态管理。", 12),
        ("3. HuggingFace社区模型数量突破50万，Transformers库月下载量超过1亿次。", 12),
        ("", 12),
        ("六、产业应用", 16),
        ("", 12),
        ("1. 金融行业：智能投研报告生成、合同审查、风控知识库问答。", 12),
        ("2. 医疗行业：病历结构化、医学文献检索、辅助诊断建议。", 12),
        ("3. 制造业：设备故障诊断知识库、工艺参数优化建议、安全规程查询。", 12),
        ("", 12),
        ("七、总结与展望", 16),
        ("", 12),
        ("2025年AI技术将继续向低成本、高效率方向演进。小模型+大上下文+Agent工具调用", 12),
        ("将成为企业级AI应用的主流架构。文档解析、RAG知识库和Agent工具调用三者", 12),
        ("的结合将赋能各行各业实现智能化升级。", 12),
    ]

    for text, fontsize in lines:
        y += 22 if fontsize == 16 else 18
        page.insert_text((72, y), text, fontsize=fontsize)

    path = OUT / "test_report.pdf"
    doc.save(str(path))
    doc.close()
    print(f"  ✓ {path}")


def generate_docx():
    """Create a Word document."""
    from docx import Document

    doc = Document()
    doc.add_heading("项目测试计划", level=1)
    doc.add_paragraph("版本: 1.0 | 日期: 2025-05-07 | 作者: 测试团队")

    doc.add_heading("一、测试范围", level=2)
    doc.add_paragraph("本次测试覆盖AI Assistant平台的以下核心功能模块：")
    items = ["文档解析（PDF/Word/PPT）", "RAG知识库问答", "Agent工具调用"]
    for item in items:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_heading("二、测试用例", level=2)
    doc.add_paragraph("TC-001: PDF文档上传并检索。预期：文档内容可被检索并返回正确答案。")
    doc.add_paragraph("TC-002: Word文档解析入库。预期：Word内容正确入库。")

    path = OUT / "test_manual.docx"
    doc.save(str(path))
    print(f"  ✓ {path}")


def generate_pptx():
    """Create a PowerPoint file."""
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000

    slides = [
        ("AI Assistant 平台介绍", ["文档解析 | RAG知识库 | Agent工具调用"]),
        ("核心能力", ["支持PDF/Word/PPT/图片解析", "基于bge-large-zh的中文语义检索", "ReAct Agent自动工具调用"]),
        ("技术栈", ["DeepSeek V4 Pro | bge-large-zh-v1.5 | pgvector", "LlamaIndex + FastAPI"]),
    ]

    for title, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        body.text = "\n".join(bullets)

    path = OUT / "test_slides.pptx"
    prs.save(str(path))
    print(f"  ✓ {path}")


def main():
    print("生成测试数据...\n")
    generate_pdf()
    generate_docx()
    generate_pptx()
    print(f"\n全部测试文件已生成到 {OUT.absolute()}")


if __name__ == "__main__":
    main()
