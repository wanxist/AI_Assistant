"""Office document parser — .docx and .pptx via python-docx and python-pptx."""

import logging
from pathlib import Path

from src.parsing.loader import ParsedDocument

logger = logging.getLogger(__name__)


class OfficeParser:
    """Parse Word (.docx) and PowerPoint (.pptx) files."""

    def parse(self, file_path: str) -> list[ParsedDocument]:
        ext = Path(file_path).suffix.lower()

        if ext == ".docx":
            return self._parse_docx(file_path)
        elif ext == ".pptx":
            return self._parse_pptx(file_path)
        else:
            raise ValueError(f"OfficeParser does not support: {ext}")

    def _parse_docx(self, file_path: str) -> list[ParsedDocument]:
        from docx import Document

        doc = Document(file_path)

        # Body paragraphs
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # Table content — each row joined as "cell1 | cell2 | ..."
        table_lines = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_lines.append(" | ".join(cells))

        # Text box content (<w:txbxContent>) — not exposed by python-docx API
        nsmap = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        txbx_blocks = doc.element.body.findall(".//w:txbxContent", nsmap)
        txbx_texts = []
        for txbx in txbx_blocks:
            texts = [t.text or "" for t in txbx.findall(".//w:t", nsmap)]
            line = "".join(texts).strip()
            if line:
                txbx_texts.append(line)

        content = "\n\n".join(paragraphs + table_lines + txbx_texts)

        return [ParsedDocument(
            file_path=file_path,
            file_type="docx",
            content=content,
            metadata={
                "source": Path(file_path).stem,
                "paragraphs": len(paragraphs),
                "tables": len(doc.tables),
                "textboxes": len(txbx_texts),
            },
            parser_used="python-docx",
        )]

    def _parse_pptx(self, file_path: str) -> list[ParsedDocument]:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
            if slide_lines:
                slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_lines))

        content = "\n\n".join(slides_text)

        return [ParsedDocument(
            file_path=file_path,
            file_type="pptx",
            content=content,
            metadata={"source": Path(file_path).stem, "slides": len(prs.slides)},
            parser_used="python-pptx",
        )]
